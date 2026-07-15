"""build_presentation.py  (NEW — makes the progress-meeting slide deck)

Generates an accessible, editable PowerPoint summarising the HateMM reproduction
(Phase 1) — the story, the results (with an ours-vs-paper chart), honest caveats,
and next steps. Numbers come from runs/phase1/reproduction_summary.md.

Output: HateMM_Reproduction_Presentation.pptx (repo root)
Also writes the results chart to runs/phase1/figs/results_macroF1.png
"""

import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = "/home/gharem/Work/Dissertation/HateMM"
FIG_DIR = os.path.join(ROOT, "runs/phase1/figs")
os.makedirs(FIG_DIR, exist_ok=True)
OUT = os.path.join(ROOT, "HateMM_Reproduction_Presentation.pptx")

# ---------- palette (colorblind-safe pair for the 2 series; dark ink text) ----------
BLUE   = RGBColor(0x3B, 0x75, 0xAF)   # "Ours"
ORANGE = RGBColor(0xE0, 0x82, 0x14)   # "Paper"
INK    = RGBColor(0x1A, 0x1A, 0x2E)
MUTED  = RGBColor(0x5B, 0x61, 0x6B)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF8)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
HEX_BLUE, HEX_ORANGE = "#3B75AF", "#E08214"

# ---------- results chart ----------
def make_chart():
    models = ["T4\ntext", "A2\naudio", "V3\nvideo", "M1", "M2", "M3", "M4"]
    ours  = [.749, .657, .697, .742, .751, .762, .767]
    paper = [.733, .669, .733, .790, .765, .767, .756]
    x = np.arange(len(models)); w = 0.38
    fig, ax = plt.subplots(figsize=(11, 5.0), dpi=200)
    b1 = ax.bar(x - w/2, ours,  w, label="Ours",  color=HEX_BLUE)
    b2 = ax.bar(x + w/2, paper, w, label="Paper", color=HEX_ORANGE)
    ax.set_ylim(0, 0.90)
    ax.set_ylabel("macro-F1  (higher is better)", fontsize=12)
    ax.set_title("Reproduction vs Paper — macro-F1 across all 7 models", fontsize=14, weight="bold")
    ax.set_xticks(x); ax.set_xticklabels(models, fontsize=11)
    ax.legend(frameon=False, ncol=2, loc="upper center", bbox_to_anchor=(0.5, 1.02), fontsize=11)
    for bars in (b1, b2):
        for r in bars:
            ax.annotate(f"{r.get_height():.2f}", (r.get_x()+r.get_width()/2, r.get_height()),
                        ha="center", va="bottom", fontsize=8, color="#333333")
    ax.yaxis.grid(True, color="#E6E6E6"); ax.set_axisbelow(True)
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    fig.tight_layout()
    path = os.path.join(FIG_DIR, "results_macroF1.png")
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path

# ---------- pptx helpers ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    return prs.slides.add_slide(BLANK)

def box(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame

def set_text(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Calibri"
    return p

def title_bar(s, text, kicker=None):
    # thin accent bar + title
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.14), Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tf = box(s, 0.9, 0.45, 11.8, 1.1)
    if kicker:
        set_text(tf, kicker.upper(), 12, MUTED, bold=True)
        p = tf.add_paragraph(); r = p.add_run(); r.text = text
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = INK; r.font.name = "Calibri"
    else:
        set_text(tf, text, 30, INK, bold=True)

def bullets(s, items, left=1.0, top=1.9, width=11.4, height=4.8, size=19, gap=10):
    tf = box(s, left, top, width, height); tf.word_wrap = True
    first = True
    for it in items:
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = lvl
        bullet = "—  " if lvl else "•  "
        r = p.add_run(); r.text = bullet + it
        r.font.size = Pt(size - (2 if lvl else 0))
        r.font.color.rgb = INK if lvl == 0 else MUTED
        r.font.name = "Calibri"

def footer(s, n):
    tf = box(s, 0.6, 7.02, 12.1, 0.4)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Video Bias Detection · Reproducing HateMM"
    r.font.size = Pt(9); r.font.color.rgb = MUTED; r.font.name = "Calibri"
    tf2 = box(s, 12.2, 7.02, 0.9, 0.4)
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = str(n)
    r2.font.size = Pt(9); r2.font.color.rgb = MUTED; r2.font.name = "Calibri"

def chip(s, l, t, w, h, title, body, color):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = LIGHT; sh.line.color.rgb = color; sh.line.width = Pt(2)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = body
    r2.font.size = Pt(13); r2.font.color.rgb = INK; r2.font.name = "Calibri"

_n = [0]
def page(): _n[0] += 1; return _n[0]

# ================= SLIDE 1 — Title =================
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(2.6))
bg.fill.solid(); bg.fill.fore_color.rgb = BLUE; bg.line.fill.background()
tf = box(s, 0.9, 0.8, 11.5, 1.4)
set_text(tf, "Video Bias Detection", 46, WHITE, bold=True)
tf = box(s, 0.9, 1.75, 11.5, 0.7)
set_text(tf, "Reproducing HateMM — Multimodal Hate-Video Detection", 22, RGBColor(0xE7,0xEE,0xF6))
tf = box(s, 0.9, 3.1, 11.5, 2.5)
set_text(tf, "MSc Dissertation · Trinity College Dublin", 18, INK, bold=True)
for line in ["Phase 1 (Reproduction) — Progress Meeting",
             "Presented by: [Your Name]",
             "July 2026"]:
    p = tf.add_paragraph(); r = p.add_run(); r.text = line
    r.font.size = Pt(16); r.font.color.rgb = MUTED; r.font.name = "Calibri"; p.space_after = Pt(4)

# ================= SLIDE 2 — Problem =================
s = slide(); title_bar(s, "Hate speech has moved to video", "The problem")
bullets(s, [
    "Video is now the dominant online format — and some of it attacks people by race, religion, or identity.",
    "Most hate-speech research is text-only. Video is far harder to screen automatically.",
    "Why? A video hides its meaning across three channels at once:",
    ("the words spoken,  the tone of voice,  and the on-screen imagery.", 1),
    "We need systems that can watch, listen, AND read — together.",
])
footer(s, page())

# ================= SLIDE 3 — 3 senses =================
s = slide(); title_bar(s, "Use three “senses” together", "The key idea")
chip(s, 0.9, 2.1, 3.7, 2.6, "Text", "What is said —\nthe speech transcript\n(the words).", BLUE)
chip(s, 4.85, 2.1, 3.7, 2.6, "Audio", "How it sounds —\ntone, shouting,\naggression.", ORANGE)
chip(s, 8.8, 2.1, 3.7, 2.6, "Video", "What is shown —\non-screen imagery,\nsymbols, targets.", GREEN)
tf = box(s, 0.9, 5.0, 11.6, 1.2)
set_text(tf, "The paper's central claim:  combining all three detects hate better than any single one alone.",
         18, INK, bold=True)
footer(s, page())

# ================= SLIDE 4 — Method =================
s = slide(); title_bar(s, "Reproduce, then improve", "My approach")
bullets(s, [
    "Phase 1 — Reproduce a published baseline (HateMM, ICWSM 2023) to build a solid, trusted foundation.   ✅ DONE",
    "Phase 2 & 3 — Add my own original contributions on top (coming next).",
    "Why reproduce first? It proves the foundation works before I build new ideas on it — and independently confirms the paper's results.",
    "Today's meeting covers the Phase 1 results.",
])
footer(s, page())

# ================= SLIDE 5 — Dataset =================
s = slide(); title_bar(s, "The dataset: HateMM", "What we work with")
bullets(s, [
    "1,083 real videos from BitChute (a low-moderation video platform).",
    "431 hateful  ·  652 non-hateful  —  each human-labelled.",
    "Two bonus labels I'll use later for my contributions:",
    ("which group is targeted (e.g., Blacks, Jews, …)", 1),
    ("the exact seconds of the video that are hateful", 1),
    "~43 hours of video, ~144,000 frames in total.",
])
footer(s, page())

# ================= SLIDE 6 — Why reproduction was hard =================
s = slide(); title_bar(s, "“Reproducing” was real engineering", "Not just re-running code")
bullets(s, [
    "The authors shared code — but it was not runnable as delivered.",
    "I had to write 3 missing scripts they never released:",
    ("data-splitting, audio extraction, and speech-to-text transcription", 1),
    "I fixed many bugs: broken file paths, syntax errors, outdated library calls — and rewrote one entirely broken feature script.",
    "The speech-to-text step alone ran for ~13 hours.",
    "Outcome: a working, documented, fully reproducible pipeline.",
])
footer(s, page())

# ================= SLIDE 7 — Pipeline =================
s = slide(); title_bar(s, "The pipeline, end to end", "How it works")
def pbox(l, t, w, h, text, color, fill=LIGHT, tsize=14, tcolor=INK, bold=True):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.color.rgb = color; sh.line.width = Pt(2)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(tsize); r.font.bold = bold
    r.font.color.rgb = tcolor; r.font.name = "Calibri"
def arrow(l, t, w=0.5, h=0.5):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = MUTED; a.line.fill.background()
pbox(0.7, 2.6, 2.2, 1.2, "1,083\nVideos", BLUE)
arrow(3.0, 2.9)
pbox(3.65, 2.6, 3.0, 1.2, "Turn each sense\ninto numbers\n(features)", ORANGE)
arrow(6.8, 2.9)
pbox(7.45, 2.6, 2.4, 1.2, "Train\nmodels", GREEN)
arrow(9.95, 2.9)
pbox(10.6, 2.6, 2.1, 1.2, "Compare\nto paper", BLUE)
# three-sense sub-row feeding features
for i, (lbl, col) in enumerate([("Text", BLUE), ("Audio", ORANGE), ("Video", GREEN)]):
    pbox(3.65 + i*1.02, 4.3, 0.95, 0.7, lbl, col, tsize=11)
tf = box(s, 0.7, 5.4, 12.0, 1.0)
set_text(tf, "Prep the ingredients (videos → numbers)  →  cook (train models)  →  taste-test against the original recipe (the paper).",
         15, MUTED)
footer(s, page())

# ================= SLIDE 8 — Features =================
s = slide(); title_bar(s, "Turning videos into numbers", "“Features”")
bullets(s, [
    "A computer can't read a video directly — we convert each sense into a list of numbers (its “features”).",
    "Text  →  768 numbers per video   (a language model, BERT).",
    "Audio →  40 and 1,000 numbers   (sound fingerprints, MFCC + VGG19).",
    "Video →  100 snapshots × 768 numbers   (an image model, ViT).",
    "Done for all 1,083 videos — the heavy, GPU-powered step (now complete).",
])
footer(s, page())

# ================= SLIDE 9 — The 7 models =================
s = slide(); title_bar(s, "The models we trained", "7 experiments")
bullets(s, [
    "3 single-sense models:  Text-only,  Audio-only,  Video-only.",
    "4 combined models:  different “recipes” that mix all three senses.",
    "Each is tested fairly with 5-fold cross-validation:",
    ("trained 5 times on different splits of the data, then averaged — so a score isn't a fluke.", 1),
    "We compare every model's score against the paper's reported score.",
])
footer(s, page())

# ================= SLIDE 10 — Results (chart) =================
s = slide(); title_bar(s, "Results: we reproduced the paper", "Ours vs Paper")
chart = make_chart()
s.shapes.add_picture(chart, Inches(0.7), Inches(1.7), width=Inches(9.4))
# side callouts
chip(s, 10.4, 1.8, 2.5, 1.5, "Within 5%", "All 7 models within 5% of the paper.", GREEN)
chip(s, 10.4, 3.5, 2.5, 1.5, "5 of 7", "within 2% — a close match.", BLUE)
tf = box(s, 0.7, 6.35, 9.4, 0.8)
set_text(tf, "macro-F1 = a balanced accuracy score (0–1, higher is better). Blue = ours, orange = the paper.",
         12, MUTED)
footer(s, page())

# ================= SLIDE 11 — Key result =================
s = slide(); title_bar(s, "The main finding holds", "Key result")
bullets(s, [
    "Combining senses beats any single sense — the paper's core claim, independently confirmed.",
    "Best single-sense model (Text):  0.749",
    "Best combined model (M4):  0.767   →  the combined model wins.",
], top=2.0, gap=14, size=20)
big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.6), Inches(11.3), Inches(1.5))
big.fill.solid(); big.fill.fore_color.rgb = GREEN; big.line.fill.background()
tf = big.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Takeaway:  multimodal  >  single-modality.  Reproduction successful."
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
footer(s, page())

# ================= SLIDE 12 — Caveats =================
s = slide(); title_bar(s, "Honest differences (good science)", "Where we differ")
bullets(s, [
    "The paper's best model (M1) scored lower for us — and our best combo was a different one (M4).",
    "The video-only model was our weakest match.",
    "Most likely reasons:",
    ("our auto-generated transcripts differ from the authors' (they didn't release theirs),", 1),
    ("and our data splits are regenerated, not identical to theirs.", 1),
    "Reporting these honestly is good practice — and they motivate my improvements.",
])
footer(s, page())

# ================= SLIDE 13 — Challenges =================
s = slide(); title_bar(s, "Challenges I solved", "Engineering effort")
bullets(s, [
    "Reconstructed a non-runnable codebase into a working, documented pipeline.",
    "Wrote the missing data-splitting, audio-extraction, and transcription steps.",
    "Ran a 13-hour speech-transcription job; handled videos missing audio or video so none were dropped.",
    "The VM's GPU driver dropped near the end → ran final training on CPU (identical results).",
])
footer(s, page())

# ================= SLIDE 14 — Next steps =================
s = slide(); title_bar(s, "What's next", "Contributions")
chip(s, 0.9, 2.0, 5.7, 2.9, "Contribution A — Target",
     "Also predict WHICH group a\nhate video targets\n(Blacks / Jews / Other),\nas a second task alongside\nhate / not-hate.", BLUE)
chip(s, 6.9, 2.0, 5.7, 2.9, "Contribution B — Explainability",
     "Make the model explain itself:\nwhich sense drove the decision,\nand which seconds were hateful —\nchecked against the human-marked\ntime-spans (ground truth).", ORANGE)
tf = box(s, 0.9, 5.2, 11.7, 1.2)
set_text(tf, "Near-term:  restore the GPU driver, then begin Contribution A.", 18, INK, bold=True)
footer(s, page())

# ================= SLIDE 15 — Close =================
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = BLUE; bg.line.fill.background()
tf = box(s, 1.0, 2.7, 11.3, 2.0)
set_text(tf, "Thank you — questions?", 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Phase 1 reproduction complete · moving to Contributions A & B"
r.font.size = Pt(18); r.font.color.rgb = RGBColor(0xE7,0xEE,0xF6); r.font.name = "Calibri"

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
